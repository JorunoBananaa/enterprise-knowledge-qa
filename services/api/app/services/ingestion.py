import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from app.services.ocr import recognize_text


@dataclass
class ParsedChunk:
    text: str
    locator: str
    chunk_index: int


# ── 分块 ──────────────────────────────────────────────────────────

_CHUNK_SIZE = 1000   # 每个块的字符数
_CHUNK_OVERLAP = 200


def _split_text_into_chunks(text: str, chunk_size: int = _CHUNK_SIZE, overlap: int = _CHUNK_OVERLAP) -> list[str]:
    """按句子边界将文本切分为有重叠的块。"""
    if not text.strip():
        return []

    sentences = re.split(r"(?<=[。！？.!?])\s*", text)
    chunks: list[str] = []
    current = ""

    for sentence in sentences:
        if len(current) + len(sentence) <= chunk_size:
            current += sentence
        else:
            if current:
                chunks.append(current.strip())
            # 保留重叠内容：带上前一个块的尾部
            if len(current) > overlap:
                current = current[-overlap:] + sentence
            else:
                current = sentence

    if current.strip():
        chunks.append(current.strip())

    return chunks


# ── 解析器 ───────────────────────────────────────────────────────────

def _parse_pdf(path: str) -> list[ParsedChunk]:
    """从 PDF 页面和嵌入图片中提取文本。"""
    from pypdf import PdfReader

    chunks: list[ParsedChunk] = []
    reader = PdfReader(path)
    image_text_by_page = _extract_pdf_image_text_by_page(reader)
    for page_idx, page in enumerate(reader.pages, start=1):
        page_text_parts: list[str] = []
        text = page.extract_text()
        if text:
            page_text_parts.append(text)
        page_text_parts.extend(image_text_by_page.get(page_idx - 1, []))

        full_page_text = "\n".join(page_text_parts)
        if not full_page_text.strip():
            continue
        page_chunks = _split_text_into_chunks(full_page_text)
        for i, chunk_text in enumerate(page_chunks):
            chunks.append(ParsedChunk(
                text=chunk_text,
                locator=f"page {page_idx}",
                chunk_index=len(chunks),
            ))
    return chunks


def _extract_pdf_image_text_by_page(reader) -> dict[int, list[str]]:
    image_text_by_page: dict[int, list[str]] = {}
    for page_idx, page in enumerate(reader.pages):
        page_image_texts: list[str] = []
        for image in getattr(page, "images", []):
            text = _recognize_image_blob(
                getattr(image, "data", b""),
                Path(getattr(image, "name", "")).suffix or ".img",
            )
            if text:
                page_image_texts.append(text)
        if page_image_texts:
            image_text_by_page[page_idx] = page_image_texts
    return image_text_by_page


def _parse_docx(path: str) -> list[ParsedChunk]:
    """从 DOCX 段落和嵌入图片中提取文本。"""
    from docx import Document

    doc = Document(path)
    full_text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text_parts.extend(_extract_docx_image_text(doc))

    full_text = "\n".join(full_text_parts)
    text_chunks = _split_text_into_chunks(full_text)
    return [
        ParsedChunk(text=t, locator="DOCX document", chunk_index=i)
        for i, t in enumerate(text_chunks)
    ]


def _extract_docx_image_text(doc) -> list[str]:
    image_texts: list[str] = []
    image_parts = [
        part for part in doc.part.related_parts.values()
        if getattr(part, "content_type", "").startswith("image/")
    ]

    for image_part in image_parts:
        text = _recognize_image_blob(
            image_part.blob,
            Path(str(image_part.partname)).suffix or ".img",
        )
        if text:
            image_texts.append(text)

    return image_texts


def _parse_pptx(path: str) -> list[ParsedChunk]:
    """从 PPTX 幻灯片和嵌入图片中提取文本。"""
    from pptx import Presentation

    prs = Presentation(path)
    full_text_parts: list[str] = []
    image_text_by_slide = _extract_pptx_image_text_by_slide(prs)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        slide_text = " ".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
        )
        if slide_text:
            slide_text_parts.append(slide_text)
        slide_text_parts.extend(image_text_by_slide.get(slide_idx - 1, []))
        if slide_text_parts:
            full_text_parts.append(f"[Slide {slide_idx}] " + "\n".join(slide_text_parts))

    full_text = "\n\n".join(full_text_parts)
    text_chunks = _split_text_into_chunks(full_text)
    return [
        ParsedChunk(text=t, locator="PPTX presentation", chunk_index=i)
        for i, t in enumerate(text_chunks)
    ]


def _extract_pptx_image_text_by_slide(prs) -> dict[int, list[str]]:
    image_text_by_slide: dict[int, list[str]] = {}
    for slide_idx, slide in enumerate(prs.slides):
        slide_image_texts: list[str] = []
        for shape in slide.shapes:
            image = getattr(shape, "image", None)
            if image is None:
                continue

            suffix = f".{getattr(image, 'ext', '')}" if getattr(image, "ext", "") else ".img"
            text = _recognize_image_blob(getattr(image, "blob", b""), suffix)
            if text:
                slide_image_texts.append(text)
        if slide_image_texts:
            image_text_by_slide[slide_idx] = slide_image_texts
    return image_text_by_slide


def _parse_xlsx(path: str) -> list[ParsedChunk]:
    """从 XLSX 工作表和嵌入图片中提取文本。"""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=False, data_only=True)
    full_text_parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text: list[str] = []
        for row in ws.iter_rows(values_only=True):
            row_str = " | ".join(str(c) for c in row if c is not None)
            if row_str.strip():
                rows_text.append(row_str)
        if rows_text:
            full_text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))
    full_text_parts.extend(_extract_xlsx_image_text(wb))
    wb.close()

    full_text = "\n\n".join(full_text_parts)
    text_chunks = _split_text_into_chunks(full_text)
    return [
        ParsedChunk(text=t, locator="XLSX spreadsheet", chunk_index=i)
        for i, t in enumerate(text_chunks)
    ]


def _extract_xlsx_image_text(workbook) -> list[str]:
    """从 XLSX 工作簿中提取嵌入图片的 OCR 文本。"""
    image_texts: list[str] = []
    for worksheet in workbook.worksheets:
        for image in getattr(worksheet, "_images", []):
            image_data = image._data() if callable(getattr(image, "_data", None)) else b""
            text = _recognize_image_blob(
                image_data,
                Path(getattr(image, "path", "")).suffix or ".img",
            )
            if text:
                image_texts.append(text)
    return image_texts


def _recognize_image_blob(blob: bytes, suffix: str) -> str:
    """将图片二进制数据写入临时文件并通过 OCR 识别文本。"""
    if not blob:
        return ""

    temp_path = ""
    normalized_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=normalized_suffix) as temp_file:
            temp_file.write(blob)
            temp_path = temp_file.name

        return recognize_text(temp_path).strip()
    finally:
        if temp_path:
            Path(temp_path).unlink(missing_ok=True)


def _parse_txt(path: str) -> list[ParsedChunk]:
    """从纯文本文件中提取文本。"""
    text = Path(path).read_text(encoding="utf-8")
    text_chunks = _split_text_into_chunks(text)
    return [
        ParsedChunk(text=t, locator="text file", chunk_index=i)
        for i, t in enumerate(text_chunks)
    ]


# ── 解析器注册表 ───────────────────────────────────────────────────────

_PARSER_REGISTRY: dict[str, Callable[[str], list[ParsedChunk]]] = {
    "pdf":  _parse_pdf,
    "docx": _parse_docx,
    "doc":  _parse_docx,
    "pptx": _parse_pptx,
    "ppt":  _parse_pptx,
    "xlsx": _parse_xlsx,
    "xls":  _parse_xlsx,
    "txt":  _parse_txt,
    "md":   _parse_txt,
    "csv":  _parse_txt,
}


def parse_document(storage_path: str, file_type: str) -> list[ParsedChunk]:
    """返回包含文本和引用元数据的解析块。

    根据 file_type 扩展名路由到对应解析器。
    """
    ext = file_type.lstrip(".").lower()
    parser = _PARSER_REGISTRY.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {file_type}")

    return parser(storage_path)
